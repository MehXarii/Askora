import fitz  # pymupdf
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from config import CHUNK_SIZE, CHUNK_OVERLAP
import os

# ─── OCR Engine Paths (Windows) ──────────────────────────────
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
POPPLER_PATH = r"C:\Users\Dell\Downloads\poppler-26.02.0\Library\bin"


def is_scanned_pdf(file_path: str) -> bool:
    """Check if a PDF has little/no extractable text (likely scanned)."""
    doc = fitz.open(file_path)
    total_text = ""
    for page in doc:
        total_text += page.get_text()
    doc.close()
    return len(total_text.strip()) < 50  # almost no text = scanned


def extract_text_from_scanned_pdf(file_path: str) -> list[dict]:
    """Use OCR to extract text from scanned/image PDFs."""
    images = convert_from_path(file_path, poppler_path=POPPLER_PATH)
    pages = []

    for page_num, image in enumerate(images, start=1):
        text = pytesseract.image_to_string(image)
        if text.strip():
            pages.append({
                "page": page_num,
                "text": text,
                "source": file_path
            })

    return pages


def extract_text_from_pdf(file_path: str) -> list[dict]:
    """Extract text from PDF, automatically using OCR if it's scanned."""
    if is_scanned_pdf(file_path):
        return extract_text_from_scanned_pdf(file_path)

    doc = fitz.open(file_path)
    pages = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        if text.strip():
            pages.append({
                "page": page_num + 1,
                "text": text,
                "source": file_path
            })
    doc.close()
    return pages


def extract_text_from_docx(file_path: str) -> list[dict]:
    doc = Document(file_path)
    pages = []
    current_text = ""
    page_num = 1

    for para in doc.paragraphs:
        if para.text.strip():
            current_text += para.text + "\n"
            if len(current_text.split()) >= 500:
                pages.append({
                    "page": page_num,
                    "text": current_text.strip(),
                    "source": file_path
                })
                current_text = ""
                page_num += 1

    if current_text.strip():
        pages.append({
            "page": page_num,
            "text": current_text.strip(),
            "source": file_path
        })

    return pages


def extract_text_from_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    pages = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"

        if slide_text.strip():
            pages.append({
                "page": slide_num,
                "text": slide_text.strip(),
                "source": file_path
            })

    return pages


def extract_text(file_path: str) -> list[dict]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(file_path)
    else:
        return []


def chunk_pages(pages: list[dict]) -> list[dict]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP
    )
    chunks = []
    for page in pages:
        splits = splitter.split_text(page["text"])
        for split in splits:
            chunks.append({
                "text": split,
                "page": page["page"],
                "source": page["source"]
            })
    return chunks


def process_file(file_path: str) -> list[dict]:
    pages = extract_text(file_path)
    chunks = chunk_pages(pages)
    return chunks


# Keep backward compatibility
def process_pdf(file_path: str) -> list[dict]:
    return process_file(file_path)